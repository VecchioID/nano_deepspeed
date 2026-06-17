"""
Generate slide deck for Chapter 2: DeepSpeed ZeRO
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ----- Color Palette -----
C_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)     # Deep blue
C_SECONDARY = RGBColor(0x25, 0x63, 0xEB)   # Blue
C_ACCENT = RGBColor(0xDC, 0x26, 0x26)      # Red
C_ACCENT_GREEN = RGBColor(0x16, 0xA3, 0x4A)
C_ACCENT_ORANGE = RGBColor(0xEA, 0x58, 0x0C)
C_BG = RGBColor(0x0F, 0x17, 0x2A)          # Dark bg
C_BG2 = RGBColor(0x1E, 0x29, 0x3B)
C_TEXT = RGBColor(0xE2, 0xE8, 0xF0)
C_TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0x33, 0x48, 0x5E)
C_LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
C_CODE_BG = RGBColor(0x0F, 0x17, 0x2A)

def add_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color=C_BG2, left=0, top=0, width=None, height=None):
    if width is None: width = prs.slide_width
    if height is None: height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Noto Sans CJK SC"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, font_size=16, color=C_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4),
                  font_name="Noto Sans CJK SC"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before: p.space_before = space_before
    if space_after: p.space_after = space_after
    return p

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_shape_bg(slide, C_CODE_BG, left, top, width, height)
    tf = add_text_box(slide, left+Inches(0.2), top+Inches(0.15),
                      width-Inches(0.4), height-Inches(0.3),
                      code_text, font_size=font_size,
                      color=RGBColor(0xE2, 0xE8, 0xF0),
                      font_name="Source Code Pro").text_frame
    return shape

def add_card(slide, left, top, width, height, title, body_lines,
             title_color=C_PRIMARY, bg_color=C_BG2):
    shape = add_shape_bg(slide, bg_color, left, top, width, height)
    tf = add_text_box(slide, left+Inches(0.25), top+Inches(0.15),
                      width-Inches(0.5), Inches(0.5),
                      title, font_size=18, color=title_color, bold=True).text_frame
    y_offset = top + Inches(0.6)
    for line in body_lines:
        tb = add_text_box(slide, left+Inches(0.25), y_offset,
                          width-Inches(0.5), Inches(0.35),
                          line, font_size=14, color=C_TEXT_DIM)
        y_offset += Inches(0.3)
    return shape

def add_bullet_text(tf, text, level=0, font_size=15, color=C_TEXT, bold=False, space_before=Pt(2), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Noto Sans CJK SC"
    p.level = level
    if space_before: p.space_before = space_before
    if space_after: p.space_after = space_after
    return p

def add_accent_bar(slide, left, top, width=Inches(0.08), height=None):
    if height is None: height = Inches(0.05)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ACCENT
    shape.line.fill.background()
    return shape

def add_color_bar(slide, left, top, width, height, color=C_SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_formula_box(slide, left, top, width, height, formula, font_size=16):
    shape = add_shape_bg(slide, RGBColor(0x1E, 0x29, 0x3B), left, top, width, height)
    tf = add_text_box(slide, left+Inches(0.2), top+Inches(0.1),
                      width-Inches(0.4), height-Inches(0.2),
                      formula, font_size=font_size, color=RGBColor(0xFB, 0xF2, 0x4F),
                      font_name="Source Code Pro").text_frame
    return shape

def page_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=C_TEXT_DIM,
                 alignment=PP_ALIGN.RIGHT)

def section_title(slide, num, title, subtitle=""):
    add_accent_bar(slide, Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.5))
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6),
                 f"0{num}" if num < 10 else str(num), font_size=14, color=C_SECONDARY, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.7),
                 title, font_size=32, color=C_WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=C_TEXT_DIM)
    # Separator
    add_color_bar(slide, Inches(0.8), Inches(1.9), Inches(2), Inches(0.04), C_SECONDARY)

TOTAL_SLIDES = 25

# ========================
# SLIDE 1: Title
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative elements
add_color_bar(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C_SECONDARY)
add_shape_bg(slide, RGBColor(0x1E, 0x29, 0x3B), Inches(0.3), Inches(0), Inches(0.08), Inches(7.5))

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "DeepSpeed ZeRO 优化详解", font_size=44, color=C_WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.6),
             "Zero Redundancy Optimizer — 零冗余优化器", font_size=22, color=C_TEXT_DIM)

add_color_bar(slide, Inches(1.5), Inches(3.5), Inches(3), Inches(0.04), C_SECONDARY)

add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10), Inches(0.5),
             "大模型训练框架系统学习 · 第 02 讲", font_size=18, color=C_TEXT)
add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.5),
             "nano_training_framework", font_size=14, color=C_TEXT_DIM)

page_number(slide, 1, TOTAL_SLIDES)

# ========================
# SLIDE 2: 目录
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 0, "目录", "CONTENTS")

items = [
    ("01", "为什么需要 ZeRO？", "DDP 的显存冗余问题"),
    ("02", "ZeRO 核心思想", "分片: 存 1/N, 用时收集"),
    ("03", "ZeRO-1: 优化器状态分片", "Adam 的 momentum + variance 分片"),
    ("04", "ZeRO-2: 梯度分片", "Reduce-Scatter 替代 All-Reduce"),
    ("05", "ZeRO-3: 参数分片", "动态参数物化, 时间换空间"),
    ("06", "通信优化", "计算-通信重叠 + Bucket 合并"),
    ("07", "ZeRO-Offload / Infinity", "CPU / NVMe 三级存储"),
    ("08", "ZeRO vs TP", "何时用 ZeRO? 何时用 TP?"),
    ("09", "训练流程", "ZeRO-3 完整伪代码"),
]

y = Inches(2.2)
for num, title, desc in items:
    add_text_box(slide, Inches(1), y, Inches(0.8), Inches(0.4),
                 num, font_size=18, color=C_SECONDARY, bold=True)
    add_text_box(slide, Inches(1.8), y, Inches(5), Inches(0.4),
                 title, font_size=18, color=C_WHITE, bold=True)
    add_text_box(slide, Inches(7.5), y, Inches(4.5), Inches(0.4),
                 desc, font_size=14, color=C_TEXT_DIM)
    y += Inches(0.55)

page_number(slide, 2, TOTAL_SLIDES)

# ========================
# SLIDE 3: 为什么需要 ZeRO
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 1, "为什么需要 ZeRO？", "DDP 中的显存冗余")

# Left: DDP memory breakdown
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.5),
             "标准 DDP (N 张卡), Adam 优化器", font_size=18, color=C_WHITE, bold=True)

memory_items = [
    ("FP16 参数", "2Ψ", "每卡一份, 共 N 份"),
    ("FP16 梯度", "2Ψ", "每卡一份, 共 N 份"),
    ("FP32 主权重", "4Ψ", "每卡一份, 共 N 份"),
    ("Adam Momentum", "4Ψ", "每卡一份, 共 N 份"),
    ("Adam Variance", "4Ψ", "每卡一份, 共 N 份"),
    ("合计 (模型状态)", "16Ψ", "被复制了 N 次!"),
]
y = Inches(2.9)
for name, size, note in memory_items:
    is_total = "合计" in name
    c = C_ACCENT if is_total else C_TEXT
    sz = 16 if is_total else 14
    add_text_box(slide, Inches(1.2), y, Inches(3), Inches(0.35), name, font_size=sz, color=c)
    add_text_box(slide, Inches(4.5), y, Inches(1.5), Inches(0.35), size, font_size=sz, color=C_SECONDARY, bold=True)
    add_text_box(slide, Inches(6), y, Inches(2), Inches(0.35), note, font_size=sz, color=C_TEXT_DIM)
    y += Inches(0.38)

# Right: the insight box
box = add_shape_bg(slide, RGBColor(0x1E, 0x3A, 0x5F), Inches(8.5), Inches(2.3), Inches(4.2), Inches(4.0))
tf = add_text_box(slide, Inches(8.8), Inches(2.5), Inches(3.8), Inches(0.8),
                  "💡 关键洞察", font_size=20, color=RGBColor(0x60, 0xA5, 0xFA), bold=True).text_frame
add_bullet_text(tf, "模型参数 + 优化器状态占显存大头", font_size=15, color=C_TEXT)
add_bullet_text(tf, "在 N 卡 DDP 中它们被复制了 N 份", font_size=15, color=C_TEXT)
add_bullet_text(tf, "如果消除冗余 → 每卡显存降 N 倍", font_size=15, color=C_ACCENT_GREEN, bold=True)
add_bullet_text(tf, "这就是 ZeRO 的出发点", font_size=15, color=C_TEXT)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.5),
             "例: LLaMA-65B, Ψ=65×10⁹, N=64 卡 → 每卡 16Ψ = 1040GB → 分片后 16Ψ/64 ≈ 16GB",
             font_size=13, color=C_TEXT_DIM)

page_number(slide, 3, TOTAL_SLIDES)

# ========================
# SLIDE 4: ZeRO 核心思想
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 2, "ZeRO 核心思想", "Zero Redundancy Optimizer")

add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.5),
             "将模型状态分片到所有数据并行进程中，每个进程只持有 1/N，在需要时通过通信集体获取",
             font_size=18, color=C_WHITE)

# Two columns: Naive DDP vs ZeRO-3
add_text_box(slide, Inches(1), Inches(3.2), Inches(4), Inches(0.4),
             "Naive DDP", font_size=20, color=RGBColor(0xF8, 0x71, 0x71), bold=True)
add_text_box(slide, Inches(6.5), Inches(3.2), Inches(4), Inches(0.4),
             "ZeRO-3", font_size=20, color=RGBColor(0x34, 0xD3, 0x99), bold=True)

# DDP box
ddp_text = """┌──────────┐┌──────────┐┌──────────┐
│  P₀~₂   ││  P₀~₂   ││  P₀~₂   │
│  G₀~₂   ││  G₀~₂   ││  G₀~₂   │
│  O₀~₂   ││  O₀~₂   ││  O₀~₂   │
└──────────┘└──────────┘└──────────┘
   GPU 0      GPU 1      GPU 2
   全部复制                     """
add_code_block(slide, Inches(1), Inches(3.7), Inches(4.5), Inches(2.2), ddp_text, font_size=12)

# ZeRO box
zero_text = """┌──────────┐┌──────────┐┌──────────┐
│   P₀    ││   P₁    ││   P₂    │
│   G₀    ││   G₁    ││   G₂    │
│   O₀    ││   O₁    ││   O₂    │
└──────────┘└──────────┘└──────────┘
   GPU 0      GPU 1      GPU 2
   各持1/3                    """
add_code_block(slide, Inches(6.5), Inches(3.7), Inches(4.5), Inches(2.2), zero_text, font_size=12)

# Bottom: summary
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             "ZeRO = ZeRO Redundancy Optimizer — 消除冗余, 节省 N 倍显存, 仅增加少量通信",
             font_size=16, color=C_TEXT_DIM)

page_number(slide, 4, TOTAL_SLIDES)

# ========================
# SLIDE 5: ZeRO 三阶段总览
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 3, "ZeRO 三阶段总览", "Stage 1 → 2 → 3, 逐步消除冗余")

# Three cards
cards = [
    ("ZeRO-1", "优化器状态分片", C_SECONDARY, [
        "仅分片: Adam momentum,",
        "variance, FP32 主权重",
        "梯度不参与分片",
        "前向/反向: 零额外通信",
        "→ 节省约 4× 优化器显存",
    ]),
    ("ZeRO-2", "梯度分片", C_ACCENT_GREEN, [
        "分片: 优化器 + 梯度",
        "Reduce-Scatter 替代",
        "All-Reduce",
        "通信量减半 (2Ψ→Ψ)",
        "→ 节省约 8× 显存",
    ]),
    ("ZeRO-3", "参数分片", C_ACCENT_ORANGE, [
        "分片: 全部模型状态",
        "动态参数物化",
        "All-Gather + Reduce-Scatter",
        "通信量 3Ψ",
        "→ 节省近 N× 显存",
    ]),
]

for i, (title, subtitle, color, lines) in enumerate(cards):
    left = Inches(0.8 + i * 4.1)
    top = Inches(2.3)
    width = Inches(3.8)
    height = Inches(4.5)

    shape = add_shape_bg(slide, C_BG2, left, top, width, height)

    # Color accent bar on top
    add_color_bar(slide, left, top, width, Inches(0.06), color)

    add_text_box(slide, left+Inches(0.3), top+Inches(0.3), width-Inches(0.6), Inches(0.5),
                 title, font_size=24, color=color, bold=True)
    add_text_box(slide, left+Inches(0.3), top+Inches(0.8), width-Inches(0.6), Inches(0.4),
                 subtitle, font_size=15, color=C_TEXT_DIM)

    # Separator
    add_color_bar(slide, left+Inches(0.3), top+Inches(1.3), Inches(1.5), Inches(0.02), color)

    y = top + Inches(1.6)
    for line in lines:
        add_text_box(slide, left+Inches(0.3), y, width-Inches(0.6), Inches(0.35),
                     line, font_size=13, color=C_TEXT)
        y += Inches(0.35)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
             "三阶段可叠加使用: ZeRO-3 = ZeRO-1 + ZeRO-2 + 参数分片",
             font_size=15, color=C_TEXT_DIM, alignment=PP_ALIGN.CENTER)

page_number(slide, 5, TOTAL_SLIDES)

# ========================
# SLIDE 6: ZeRO-1 Detailed
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 4, "ZeRO-1: 优化器状态分片", "每个进程只维护 1/N 的优化器状态")

# Memory before/after
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.4),
             "显存对比 (Ψ 参数, N 卡, Adam)", font_size=18, color=C_WHITE, bold=True)

# Before
add_text_box(slide, Inches(1), Inches(2.9), Inches(5), Inches(0.4),
             "分片前:  16Ψ  =  2Ψ(参数) + 2Ψ(梯度) + 12Ψ(优化器状态)",
             font_size=14, color=RGBColor(0xF8, 0x71, 0x71))

# After
add_text_box(slide, Inches(1), Inches(3.4), Inches(5.5), Inches(0.4),
             "分片后:  4Ψ + 12Ψ/N  =  2Ψ(参数) + 2Ψ(梯度) + 12Ψ/N(优化器状态)",
             font_size=14, color=RGBColor(0x34, 0xD3, 0x99))

# Formula box
add_formula_box(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.8),
                "M_ZeRO1 = 2Ψ + 2Ψ + 12Ψ/N = 4Ψ + 12Ψ/N", font_size=16)

# Communication on the right
box = add_shape_bg(slide, C_BG2, Inches(7), Inches(2.3), Inches(5.5), Inches(3.0))
tf = add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                  "通信行为", font_size=18, color=C_SECONDARY, bold=True).text_frame
add_bullet_text(tf, "前向传播: 零额外通信", font_size=15, color=C_TEXT)
add_bullet_text(tf, "反向传播: 零额外通信 (和 DDP 一样的 All-Reduce)", font_size=15, color=C_TEXT)
add_bullet_text(tf, "优化器更新: 每卡只更新自己分片部分 → 无需通信", font_size=15, color=C_TEXT)
add_bullet_text(tf, "更新后: 需要一次 All-Gather 同步完整参数 (仅这一步)", font_size=15, color=C_TEXT)

# Key insight
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "💡 ZeRO-1 在前向/反向阶段零额外通信 → 性能损失几乎为零",
             font_size=16, color=RGBColor(0x60, 0xA5, 0xFA))

page_number(slide, 6, TOTAL_SLIDES)

# ========================
# SLIDE 7: ZeRO-2
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 5, "ZeRO-2: 梯度分片", "Reduce-Scatter 替代 All-Reduce")

# Left: communication comparison
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(6), Inches(0.4),
             "All-Reduce  vs  Reduce-Scatter", font_size=18, color=C_WHITE, bold=True)

# All-Reduce
add_text_box(slide, Inches(1), Inches(2.9), Inches(2), Inches(0.4),
             "All-Reduce:", font_size=15, color=RGBColor(0xF8, 0x71, 0x71))
add_code_block(slide, Inches(1), Inches(3.3), Inches(5), Inches(1.0),
               "① 各卡求和 (sum)\n② 广播结果到所有卡\n通信量: 2Ψ  (sum + broadcast)", font_size=12)

# Reduce-Scatter
add_text_box(slide, Inches(1), Inches(4.6), Inches(2.5), Inches(0.4),
             "Reduce-Scatter:", font_size=15, color=RGBColor(0x34, 0xD3, 0x99))
add_code_block(slide, Inches(1), Inches(5.0), Inches(5.5), Inches(1.2),
               "① 各卡对梯度求和 (reduce)\n② 结果按 rank 切片分发 (scatter)\n   卡0拿G₀, 卡1拿G₁, ...\n通信量: Ψ  (reduce only, 减半!)", font_size=12)

# Right: memory
box = add_shape_bg(slide, C_BG2, Inches(7), Inches(2.3), Inches(5.5), Inches(3.5))
tf = add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.4),
                  "显存收益", font_size=18, color=C_SECONDARY, bold=True).text_frame
add_bullet_text(tf, "分片前: 2Ψ(参数)+2Ψ(梯度)+12Ψ(优化器)", font_size=14, color=C_TEXT)
add_bullet_text(tf, "分片后: 2Ψ + 2Ψ/N + 12Ψ/N", font_size=14, color=RGBColor(0x34, 0xD3, 0x99), bold=True)
add_bullet_text(tf, "", font_size=8, color=C_TEXT)
add_bullet_text(tf, f"例: Ψ=7B, N=8 → 梯度从 14GB → 1.75GB", font_size=13, color=C_TEXT_DIM)
add_bullet_text(tf, f"优化器从 84GB → 10.5GB", font_size=13, color=C_TEXT_DIM)
add_bullet_text(tf, "", font_size=8, color=C_TEXT)
add_bullet_text(tf, "💡 计算-通信可重叠:", font_size=14, color=RGBColor(0x60, 0xA5, 0xFA))
add_bullet_text(tf, "逐层 Reduce-Scatter, 算下一层的同时", font_size=13, color=C_TEXT_DIM)
add_bullet_text(tf, "通信上一层的梯度", font_size=13, color=C_TEXT_DIM)

page_number(slide, 7, TOTAL_SLIDES)

# ========================
# SLIDE 8: ZeRO-3
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 6, "ZeRO-3: 参数分片", "动态参数物化 (Dynamic Parameter Materialization)")

# Core mechanism diagram
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.5),
             "核心流程: 计算某层时才收集该层参数 → 计算 → 立即释放",
             font_size=18, color=C_WHITE)

# Forward process
add_text_box(slide, Inches(1), Inches(3.0), Inches(4), Inches(0.4),
             "前向传播 (Layer i):", font_size=16, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
add_code_block(slide, Inches(1), Inches(3.5), Inches(5.5), Inches(1.5),
               "① All-Gather: 收集 Layer_i 的完整参数\n   (从所有 N 个 rank)\n"
               "② 计算前向:  Layer_i(input)\n"
               "③ 释放: 非本 rank 的参数分片\n   (只保留自己应持有的 1/N)", font_size=12)

# Backward process
add_text_box(slide, Inches(7), Inches(3.0), Inches(5), Inches(0.4),
             "反向传播 (Layer i):", font_size=16, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
add_code_block(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(1.5),
               "① All-Gather: 重新收集 Layer_i 参数\n"
               "② 计算梯度:  Layer_i.backward()\n"
               "③ Reduce-Scatter: 梯度求和并分发\n"
               "④ 释放: 非本 rank 参数和梯度", font_size=12)

# Bottom: trade-off
add_accent_bar(slide, Inches(0.8), Inches(5.5), Inches(0.08), Inches(1.2))
tf = add_text_box(slide, Inches(1.1), Inches(5.5), Inches(10), Inches(0.5),
                  "时间换空间", font_size=18, color=C_ACCENT, bold=True).text_frame
add_bullet_text(tf, "显存从 16Ψ → 16Ψ/N (减少近 N 倍)", font_size=15, color=C_TEXT)
add_bullet_text(tf, "通信从 2Ψ (DDP) → 3Ψ (ZeRO-3: 前向 All-Gather + 反向 All-Gather + Reduce-Scatter)", font_size=15, color=C_TEXT)

page_number(slide, 8, TOTAL_SLIDES)

# ========================
# SLIDE 9: 三阶段显存对比
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 7, "三阶段显存对比", "N=64 卡, Adam 优化器, 数值归一化")

# Memory breakdown visualization
labels = ["Naive DDP", "ZeRO-1", "ZeRO-2", "ZeRO-3"]
mem_params = [24, 20.2, 16.1, 0.38]  # units
mem_breaks = [
    [("参数", 4, C_ACCENT), ("梯度", 4, RGBColor(0x60, 0xA5, 0xFA)), ("优化器", 16, C_SECONDARY)],
    [("参数", 4, C_ACCENT), ("梯度", 4, RGBColor(0x60, 0xA5, 0xFA)), ("优化器", 12.25, C_SECONDARY)],
    [("参数", 4, C_ACCENT), ("梯度", 0.0625, RGBColor(0x60, 0xA5, 0xFA)), ("优化器", 12.05, C_SECONDARY)],
    [("参数", 0.06, C_ACCENT), ("梯度", 0.06, RGBColor(0x60, 0xA5, 0xFA)), ("优化器", 0.19, C_SECONDARY)],
]

bar_left = Inches(4)
bar_top_start = Inches(2.3)
bar_width = Inches(0.8)
bar_max_height = Inches(3.8)
max_val = 24

for i, (label, val, breaks) in enumerate(zip(labels, mem_params, mem_breaks)):
    y = bar_top_start + Inches(4.0) * (1 - val / max_val)

    # label above bar
    add_text_box(slide, Inches(0.5), Inches(2.3 + i*1.0), Inches(3.2), Inches(0.4),
                 label, font_size=17, color=C_WHITE, bold=True)

    # Stacked bar
    y_pos = bar_top_start + Inches(4.0) * (1 - val / max_val)
    for br_name, br_val, br_color in breaks:
        br_h = Inches(4.0) * br_val / max_val
        if br_h < Inches(0.06):
            br_h = Inches(0.06)
        add_color_bar(slide, bar_left, y_pos, bar_width, br_h, br_color)
        y_pos += br_h

    # Value text
    add_text_box(slide, bar_left + bar_width + Inches(0.1), bar_top_start + Inches(4.0)*(1-val/max_val),
                 Inches(1.5), Inches(0.4),
                 f"{val:.1f} 单位", font_size=14, color=C_WHITE, bold=True)
    if val < 1:
        add_text_box(slide, bar_left + bar_width + Inches(0.1), bar_top_start + Inches(4.0)*(1-val/max_val)+Inches(0.3),
                     Inches(1.5), Inches(0.3),
                     f"(-{int((1-val/24)*100)}%)", font_size=12, color=RGBColor(0x34, 0xD3, 0x99))

# Legend
legend_y = Inches(6.8)
colors = [C_ACCENT, RGBColor(0x60, 0xA5, 0xFA), C_SECONDARY]
lbls = ["参数", "梯度", "优化器"]
for i, (c, l) in enumerate(zip(colors, lbls)):
    add_color_bar(slide, Inches(4 + i*2.5), legend_y, Inches(0.3), Inches(0.2), c)
    add_text_box(slide, Inches(4.4 + i*2.5), legend_y - Inches(0.03), Inches(1.5), Inches(0.3),
                 l, font_size=12, color=C_TEXT)

page_number(slide, 9, TOTAL_SLIDES)

# ========================
# SLIDE 10: 三阶段对比表
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 8, "三阶段详细对比")

# Table
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

rows, cols = 7, 4
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5))
table = table_shape.table

headers = ["维度", "ZeRO-1", "ZeRO-2", "ZeRO-3"]
data = [
    ["分片内容", "Optimizer\n状态", "Optimizer\n+ 梯度", "Optimizer\n+ 梯度 + 参数"],
    ["每卡显存\n(Ψ=N=64)", "4Ψ + 12Ψ/N\n≈ 4.19Ψ", "2Ψ + 14Ψ/N\n≈ 2.22Ψ", "16Ψ/N\n≈ 0.25Ψ"],
    ["显存节省\n(相对Naive)", "~4×", "~8×", "~64×"],
    ["额外通信量", "无", "无 (改Reduce-\nScatter,更优)", "3Ψ/step\n(2 All-Gather\n+ Reduce-Scatter)"],
    ["计算通信\n重叠", "不易", "梯度阶段\n可重叠", "前向/反向\n均可重叠"],
    ["适用场景", "小规模扩展\n(<32卡)", "中等规模\n(32-128卡)", "超大规模\n(128卡+)"],
]

# Style headers
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRIMARY

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT if j > 0 else C_WHITE
        p.font.bold = (j == 0)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        if j == 0:
            cell.fill.fore_color.rgb = C_BG2
        else:
            cell.fill.fore_color.rgb = C_BG2 if i % 2 == 0 else RGBColor(0x16, 0x20, 0x30)

page_number(slide, 10, TOTAL_SLIDES)

# ========================
# SLIDE 11: 通信优化 - Overlap
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 9, "通信优化: 计算-通信重叠", "将通信隐藏在计算背后")

# Timeline diagram
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.4),
             "ZeRO-3 重叠示意", font_size=18, color=C_WHITE, bold=True)

# Gantt-like chart
y_base = Inches(2.9)
bar_h = Inches(0.4)

# Forward
add_text_box(slide, Inches(0.3), y_base, Inches(1.5), Inches(0.4),
             "前向:", font_size=14, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)

# All-Gather W0
add_color_bar(slide, Inches(1.8), y_base, Inches(2), bar_h, C_ACCENT_ORANGE)
add_text_box(slide, Inches(1.8), y_base, Inches(2), bar_h,
             "AG W₀", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# Compute F0
add_color_bar(slide, Inches(3.8), y_base, Inches(2.5), bar_h, C_ACCENT_GREEN)
add_text_box(slide, Inches(3.8), y_base, Inches(2.5), bar_h,
             "计算 F₀", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# AG W1 (overlapping with F0)
add_color_bar(slide, Inches(5.5), y_base, Inches(2), bar_h, C_ACCENT_ORANGE)
add_text_box(slide, Inches(5.5), y_base, Inches(2), bar_h,
             "AG W₁", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.5), y_base, Inches(5), Inches(0.4),
             "All-Gather W₁ 与 计算 F₀ 同时进行!", font_size=13, color=C_TEXT_DIM)

# Backward
y_base2 = Inches(3.6)
add_text_box(slide, Inches(0.3), y_base2, Inches(1.5), Inches(0.4),
             "反向:", font_size=14, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)

# Compute BN
add_color_bar(slide, Inches(1.8), y_base2, Inches(2.5), bar_h, C_ACCENT_GREEN)
add_text_box(slide, Inches(1.8), y_base2, Inches(2.5), bar_h,
             "计算 B_N", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# RS GN
add_color_bar(slide, Inches(4.3), y_base2, Inches(2), bar_h, RGBColor(0xA7, 0x8B, 0xFA))
add_text_box(slide, Inches(4.3), y_base2, Inches(2), bar_h,
             "RS G_N", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# B N-1
add_color_bar(slide, Inches(5.8), y_base2, Inches(2.5), bar_h, C_ACCENT_GREEN)
add_text_box(slide, Inches(5.8), y_base2, Inches(2.5), bar_h,
             "计算 B_N₋₁", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(8.5), y_base2, Inches(5), Inches(0.4),
             "Reduce-Scatter G_N 与 计算 B_N₋₁ 重叠!", font_size=13, color=C_TEXT_DIM)

# Implementation details
tf = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4),
                  "实现机制", font_size=16, color=C_SECONDARY, bold=True).text_frame
add_bullet_text(tf, "异步 NCCL 通信: ncclGroupStart/End + CUDA events 同步", font_size=14, color=C_TEXT)
add_bullet_text(tf, "双 stream: 一个通信 stream, 一个计算 stream, 并行执行", font_size=14, color=C_TEXT)
add_bullet_text(tf, "预取机制: 计算当前层时, 提前发起下一层参数的 All-Gather", font_size=14, color=C_TEXT)
add_bullet_text(tf, "Bucket 合并: 多个小参数合并为一个 All-Gather, 提高带宽利用率", font_size=14, color=C_TEXT)

page_number(slide, 11, TOTAL_SLIDES)

# ========================
# SLIDE 12: ZeRO-Offload
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 10, "ZeRO-Offload", "利用 CPU 内存扩展可训练模型规模")

# Architecture diagram
# GPU section
gpu_shape = add_shape_bg(slide, C_BG2, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.0))
add_text_box(slide, Inches(1.1), Inches(2.5), Inches(2), Inches(0.4),
             "GPU HBM", font_size=16, color=C_WHITE, bold=True)

gpu_items = [
    ("参数分片 (1/N)", Inches(2.1), C_ACCENT),
    ("梯度分片 (1/N)", Inches(2.6), RGBColor(0x60, 0xA5, 0xFA)),
    ("计算 + Activation", Inches(3.1), C_ACCENT_GREEN),
]
for item, y_pos, color in gpu_items:
    add_text_box(slide, Inches(1.3), y_pos, Inches(3), Inches(0.35), item, font_size=14, color=color)
    add_color_bar(slide, Inches(4.5), y_pos, Inches(1.2), Inches(0.04), color)

# Arrow down
add_text_box(slide, Inches(3.8), Inches(3.5), Inches(1.5), Inches(0.5),
             "▼ offload", font_size=14, color=RGBColor(0xF8, 0x71, 0x71), alignment=PP_ALIGN.CENTER)

# CPU section
cpu_shape = add_shape_bg(slide, C_BG2, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.5))
add_text_box(slide, Inches(1.1), Inches(5.7), Inches(3), Inches(0.4),
             "CPU DRAM", font_size=16, color=C_WHITE, bold=True)
add_text_box(slide, Inches(1.3), Inches(6.2), Inches(4.5), Inches(0.4),
             "优化器状态 (FP32 momentum + variance)", font_size=14, color=C_SECONDARY)

# Right: notes
tf = add_text_box(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5),
                  "", font_size=18, color=C_WHITE).text_frame
p = add_paragraph(tf, "Offload 策略", font_size=18, color=C_SECONDARY, bold=True)
add_bullet_text(tf, "GPU→CPU 传输: PCIe ≈ 64 GB/s", font_size=14, color=C_TEXT)
add_bullet_text(tf, "NVLink 带宽: ~900 GB/s", font_size=14, color=C_TEXT)
add_bullet_text(tf, "相差约 14 倍!", font_size=14, color=RGBColor(0xF8, 0x71, 0x71))
add_bullet_text(tf, "", font_size=8, color=C_TEXT)
add_bullet_text(tf, "推荐: 仅 Offload 优化器状态", font_size=15, color=C_ACCENT_GREEN, bold=True)
add_bullet_text(tf, "参数和梯度保留在 GPU", font_size=14, color=C_TEXT)
add_bullet_text(tf, "→ 约 10-20% 性能损失", font_size=14, color=C_TEXT_DIM)
add_bullet_text(tf, "", font_size=8, color=C_TEXT)
add_bullet_text(tf, "不推荐: Offload 参数", font_size=14, color=RGBColor(0xF8, 0x71, 0x71))
add_bullet_text(tf, "→ 30-50% 性能损失", font_size=14, color=C_TEXT_DIM)

page_number(slide, 12, TOTAL_SLIDES)

# ========================
# SLIDE 13: ZeRO-Infinity
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 11, "ZeRO-Infinity", "NVMe 第三级存储")

# Storage hierarchy
levels = [
    ("GPU HBM", "~80 GB", "~1 μs", "~2000 GB/s", C_ACCENT),
    ("CPU DRAM", "~1 TB", "~0.1 μs", "~100 GB/s", C_ACCENT_ORANGE),
    ("NVMe SSD", "~10 TB", "~10 μs", "~7 GB/s", C_TEXT_DIM),
]

y = Inches(2.5)
for i, (name, capacity, latency, bandwidth, color) in enumerate(levels):
    left = Inches(1.5 + i * 3.8)
    shape = add_shape_bg(slide, C_BG2, left, y, Inches(3.2), Inches(2.0))
    add_text_box(slide, left+Inches(0.2), y+Inches(0.2), Inches(2.8), Inches(0.4),
                 name, font_size=18, color=color, bold=True)
    add_text_box(slide, left+Inches(0.2), y+Inches(0.7), Inches(2.8), Inches(0.3),
                 f"容量: {capacity}", font_size=13, color=C_TEXT)
    add_text_box(slide, left+Inches(0.2), y+Inches(1.0), Inches(2.8), Inches(0.3),
                 f"延迟: {latency}", font_size=13, color=C_TEXT)
    add_text_box(slide, left+Inches(0.2), y+Inches(1.3), Inches(2.8), Inches(0.3),
                 f"带宽: {bandwidth}", font_size=13, color=C_TEXT)

    if i < 2:
        add_text_box(slide, left+Inches(3.2), y+Inches(0.7), Inches(0.6), Inches(0.4),
                     "↓", font_size=20, color=C_TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
             "ZeRO-Infinity 自动根据访问频率决定 offload 层级:\n"
             "  频繁访问 (每 step) → GPU HBM\n"
             "  偶尔访问 (每 epoch) → CPU DRAM\n"
             "  很少访问 → NVMe SSD\n\n"
             "理论: 可训练无限大模型 (受限于 NVMe 带宽)",
             font_size=15, color=C_TEXT)

page_number(slide, 13, TOTAL_SLIDES)

# ========================
# SLIDE 14: ZeRO vs TP
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 12, "ZeRO vs 张量并行 (TP)", "何时用 ZeRO？何时用 TP？")

# Comparison table
rows, cols = 6, 3
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.5))
table = table_shape.table

headers_tp = ["维度", "ZeRO-3", "Tensor Parallelism"]
data_tp = [
    ["分片粒度", "逐层 (Per Layer)", "逐算子 (Per Op)"],
    ["通信量", "O(Ψ) 每层", "O(H) 每层 (H=hidden)"],
    ["消息特征", "少量大消息", "大量小消息"],
    ["跨节点性能", "好 (IB 即可)", "差 (需要 NVLink)"],
    ["扩展上限", "千卡级", "通常 ≤ 8 卡"],
]

for j, h in enumerate(headers_tp):
    cell = table.cell(0, j)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_PRIMARY

for i, row_data in enumerate(data_tp):
    for j, val in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT if j > 0 else C_WHITE
        p.font.bold = (j == 0)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        if j == 0:
            cell.fill.fore_color.rgb = C_BG2
        else:
            cell.fill.fore_color.rgb = C_BG2 if i % 2 == 0 else RGBColor(0x16, 0x20, 0x30)

# Bottom: combo recommendation
add_accent_bar(slide, Inches(0.8), Inches(6.3), Inches(0.08), Inches(0.8))
tf = add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.8),
                  "实际组合: TP (节点内, 8卡) + ZeRO-1/2 (DP组) + PP (跨节点)", font_size=16, color=C_WHITE).text_frame
add_bullet_text(tf, "TP + ZeRO-3 不推荐: 两者都频繁 All-Gather, 通信叠加 → 性能下降", font_size=14, color=C_TEXT_DIM)

page_number(slide, 14, TOTAL_SLIDES)

# ========================
# SLIDE 15: 训练伪代码
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 13, "ZeRO-3 训练流程伪代码")

code = """for each batch:
    # 前向传播
    for layer in model:
        nccl.all_gather(layer.params, from_all_ranks)  # 收集完整参数
        output = layer.forward(input)                   # 计算
        layer.params.free_except(my_rank_slice)         # 释放非本 rank 分片

    # 计算 loss
    loss = criterion(output, target)

    # 反向传播
    for layer in reversed(model):
        nccl.all_gather(layer.params, from_all_ranks)   # 重新收集参数
        layer.backward(grad_output)                     # 计算梯度
        nccl.reduce_scatter(layer.grad, to_all_ranks)   # 梯度求和并分发
        layer.params.free_except(my_rank_slice)          # 释放
        layer.grad.free_except(my_rank_slice)

    # 优化器更新 (无通信!)
    optimizer.step(my_rank_slice_params, my_rank_slice_grads)"""

add_code_block(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), code, font_size=12)

page_number(slide, 15, TOTAL_SLIDES)

# ========================
# SLIDE 16: 总结
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 14, "总结", "Key Takeaways")

takeaways = [
    ("ZeRO 的核心", "消除数据并行中的冗余存储\n每个 GPU 只存 1/N 的模型状态\n需要时通过通信集体获取 (All-Gather / Reduce-Scatter)"),
    ("三阶段递进", "ZeRO-1: 优化器分片, 零额外通信\nZeRO-2: + 梯度分片, 通信减半\nZeRO-3: + 参数分片, 显存降 N 倍"),
    ("性能关键", "计算-通信重叠是 ZeRO 的性能核心\n异步 NCCL + 双 Stream + 参数预取\nBucket 合并: 小消息 → 大消息"),
    ("Offload", "GPU → CPU → NVMe 三级存储\n仅推荐 Offload 优化器状态\n兜底策略, 性能损失 10-20%"),
]

for i, (title, desc) in enumerate(takeaways):
    left = Inches(0.8 + (i % 2) * 6.2)
    top = Inches(2.3 + (i // 2) * 2.3)
    shape = add_shape_bg(slide, C_BG2, left, top, Inches(5.8), Inches(2.0))
    add_color_bar(slide, left, top, Inches(0.06), Inches(2.0), C_SECONDARY)
    add_text_box(slide, left+Inches(0.3), top+Inches(0.2), Inches(5.2), Inches(0.4),
                 title, font_size=18, color=C_SECONDARY, bold=True)
    add_text_box(slide, left+Inches(0.3), top+Inches(0.7), Inches(5.2), Inches(1.2),
                 desc, font_size=14, color=C_TEXT)

page_number(slide, 16, TOTAL_SLIDES)

# ========================
# SLIDE 17: 思考题
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_title(slide, 15, "思考题")

questions = [
    "1. ZeRO-3 的 All-Gather 在通信量上比 DDP 的 All-Reduce 还多 50%，\n"
    "   为什么实际效果仍然很好？\n"
    "   (提示: 考虑通信重叠、逐层计算、带宽利用率)",

    "2. 为什么 ZeRO-3 和 TP 通常不一起使用？\n"
    "   如果必须一起用，应该怎么做？\n"
    "   (提示: 通信模式冲突、带宽争夺)",

    "3. ZeRO-Offload 中，为什么选择 offload optimizer states\n"
    "   而不是 gradients？\n"
    "   (提示: 访问频率、PCIe 带宽、计算-通信路径)",
]

for i, q in enumerate(questions):
    y = Inches(2.4 + i * 1.7)
    shape = add_shape_bg(slide, C_BG2, Inches(0.8), y, Inches(11.5), Inches(1.4))
    add_text_box(slide, Inches(1.1), y+Inches(0.15), Inches(11), Inches(1.2),
                 q, font_size=15, color=C_TEXT)

page_number(slide, 17, TOTAL_SLIDES)

# ========================
# SLIDE 18: Thank You
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_color_bar(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), C_SECONDARY)
add_shape_bg(slide, RGBColor(0x1E, 0x29, 0x3B), Inches(0.3), Inches(0), Inches(0.08), Inches(7.5))

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "谢谢", font_size=48, color=C_WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
             "下一讲: DeepSpeed 高级功能与实战", font_size=22, color=C_TEXT_DIM)

add_color_bar(slide, Inches(1.5), Inches(4.5), Inches(3), Inches(0.04), C_SECONDARY)

add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
             "nano_training_framework · 大模型训练框架系统学习", font_size=16, color=C_TEXT)

# Save
output_path = str(Path(__file__).parent / "pdf_output" / "02-DeepSpeed-ZeRO课件.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
